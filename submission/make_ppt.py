#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
solution.pptx 생성 스크립트 (7슬라이드)

1. 타이틀 (Score 0.6295)
2. 과제 개요 (평가 지표, 벤치마크, 제약)
3. 최종 솔루션 (QuantMod W8A8 상세)
4. 양자화 방식별 서버 결과 비교
5. 시도했으나 실패한 방법들 (110+ 모델)
6. 핵심 인사이트 4가지
7. 모델 크기 비교 + 재현 방법
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 색상 팔레트
# ============================================================
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x2D, 0x5A, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE8, 0x4D, 0x3D)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
TEXT = RGBColor(0x2C, 0x3E, 0x50)
HIGHLIGHT = RGBColor(0xE8, 0xF4, 0xFD)
SKY = RGBColor(0x85, 0xC1, 0xE9)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font_name="맑은 고딕"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return p


def add_table(slide, left, top, width, height, data, col_widths=None,
              header_color=BLUE):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if i > 0 else 12)
                p.font.name = "맑은 고딕"
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = TEXT
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


def add_divider(slide, y=0.95):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(y), Inches(12.3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=LIGHT_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# PPT 생성
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ===========================================================
# Slide 1: 타이틀
# ===========================================================
slide = prs.slides.add_slide(blank)
add_bg(slide, DARK)

add_text(slide, 1, 1.0, 11, 0.8,
         "LG Aimers Phase 2 — LLM 경량화 해커톤",
         size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 1, 2.0, 11, 1.2,
         "EXAONE-4.0-1.2B 모델 경량화",
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, 1, 3.4, 11, 0.8,
         "QuantizationModifier W8A8 + Context Window 축소",
         size=22, color=SKY, align=PP_ALIGN.CENTER)

# Score 강조
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.2), Inches(4.6), Inches(4.9), Inches(1.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Final Score: 0.6295"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "맑은 고딕"
p.alignment = PP_ALIGN.CENTER

add_text(slide, 1, 6.0, 11, 0.5,
         "모델 크기: 2.4GB → 1.4GB (42% 감소)  |  추론 속도: 1,811 → 2,435 tok/s",
         size=16, color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

add_text(slide, 1, 6.6, 11, 0.5,
         "팀원: [이름 직접 기재]  |  오프라인 참가 여부: [직접 기재]",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===========================================================
# Slide 2: 과제 개요
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 12, 0.6, "과제 개요", size=30, bold=True, color=BLUE)
add_divider(slide)

# 왼쪽: 평가 지표
tf = add_text(slide, 0.5, 1.3, 5.8, 4.5, "평가 지표", size=20, bold=True, color=BLUE)
add_para(tf, "")
add_para(tf, "Score = 0.5 × PerfNorm + 0.5 × SpeedNorm", size=16, bold=True, color=BLUE)
add_para(tf, "")
add_para(tf, "• PerfNorm = 압축모델 성능 / 베이스모델 성능", size=14)
add_para(tf, "• SpeedNorm = 1 − (압축 time/tok / 베이스 time/tok)", size=14)
add_para(tf, "")
add_para(tf, "→ 비압축 모델: Score = 0.5 (PerfNorm=1, SpeedNorm=0)", size=13, color=GRAY)
add_para(tf, "→ 성능 유지 + 속도 향상이 핵심", size=14, bold=True, color=RED)
add_para(tf, "")
add_para(tf, "평가 환경", size=20, bold=True, color=BLUE)
add_para(tf, "")
add_para(tf, "• 서버 GPU: NVIDIA L4 (24GB, Ada Lovelace)", size=14)
add_para(tf, "• 추론 엔진: vLLM 0.14.1", size=14)
add_para(tf, "• 베이스 모델: EXAONE-4.0-1.2B (2.4GB, BF16)", size=14)
add_para(tf, "• 제출 제한: 압축 10GB / 비압축 32GB", size=14)

# 오른쪽: 벤치마크
tf2 = add_text(slide, 7, 1.3, 5.8, 2.5, "평가 벤치마크", size=20, bold=True, color=BLUE)
add_para(tf2, "")
add_para(tf2, "• KMMLU-Pro (한국어 MCQA, 2,822문항)", size=14)
add_para(tf2, "• KMMLU-Redux (한국어 MCQA, 2,587문항)", size=14)
add_para(tf2, "• Ko-LongRAG (긴 문맥 QA, 600문항)", size=14)
add_para(tf2, "• KoMT-Bench (GPT-4 채점, 80문항)", size=14)

tf3 = add_text(slide, 7, 4.0, 5.8, 3.0, "핵심 제약", size=20, bold=True, color=BLUE)
add_para(tf3, "")
add_para(tf3, "• 1.2B 소형 모델 → INT4 양자화 시 품질 붕괴", size=14)
add_para(tf3, "• 로컬(RTX 5060Ti)과 서버(L4)의 최적 방식 상이", size=14)
add_para(tf3, "• L4에서 W8A8(INT8)만 실질적 속도 향상", size=14)
add_para(tf3, "• W8A16, FP8 등은 커널 비효율로 속도 이득 미미", size=14)

# ===========================================================
# Slide 3: 최종 솔루션
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 12, 0.6,
         "최종 솔루션 — QuantizationModifier W8A8",
         size=30, bold=True, color=BLUE)
add_divider(slide)

# 왼쪽: 양자화 설정 테이블
add_text(slide, 0.5, 1.2, 6, 0.4, "양자화 설정", size=18, bold=True, color=BLUE)
add_table(slide, 0.5, 1.7, 6, 3.6, [
    ["항목", "설정"],
    ["양자화 도구", "llmcompressor 0.9.0.1 (QuantizationModifier)"],
    ["양자화 스킴", "W8A8 (INT8 weight + INT8 activation)"],
    ["Weight 양자화", "per-channel, symmetric, static (minmax)"],
    ["Activation 양자화", "per-token, symmetric, dynamic"],
    ["대상 레이어", "모든 Linear"],
    ["제외 레이어", "embed_tokens, lm_head"],
    ["캘리브레이션", "MANTA-1M 2,048 샘플, seq_len=2048"],
], col_widths=[2.2, 3.8])

# 오른쪽 상단: Context Window
add_text(slide, 7, 1.2, 5.5, 0.4,
         "Context Window 최적화", size=18, bold=True, color=BLUE)
tf = add_text(slide, 7, 1.7, 5.5, 2.0, "", size=14)
add_para(tf, "max_position_embeddings:", size=14, bold=True)
add_para(tf, "65,536 → 16,384 (4배 축소)", size=16, bold=True, color=RED)
add_para(tf, "")
add_para(tf, "• KV Cache 메모리 4배 절감", size=13)
add_para(tf, "• 더 큰 배치 처리 → 추론 속도 향상", size=13)
add_para(tf, "• 벤치마크 최대 입력 ~16k → 품질 손실 없음", size=13, color=GREEN)

# 오른쪽 하단: QuantMod vs GPTQ
add_text(slide, 7, 4.0, 5.5, 0.4,
         "QuantMod vs GPTQ 비교", size=18, bold=True, color=BLUE)
tf = add_text(slide, 7, 4.5, 5.5, 2.5, "", size=13)
add_para(tf, "QuantMod (min-max): 0.6295 ← 최종 선택", size=14, bold=True, color=GREEN)
add_para(tf, "GPTQ (Hessian): 0.6208", size=14)
add_para(tf, "")
add_para(tf, "→ 역설적으로 단순한 방식이 더 높은 점수", size=13, bold=True, color=RED)
add_para(tf, "• GPTQ Hessian이 캘리브레이션 데이터에 과적합", size=12, color=GRAY)
add_para(tf, "• min-max 방식이 분포 변화에 더 robust", size=12, color=GRAY)

# ===========================================================
# Slide 4: 서버 결과 비교
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 12, 0.6,
         "양자화 방식별 서버 결과 비교", size=30, bold=True, color=BLUE)
add_divider(slide)

result_table = add_table(slide, 0.5, 1.2, 12.3, 5.0, [
    ["순위", "방법", "서버 Score", "모델 크기", "비고"],
    ["1", "QuantMod W8A8 + ctx16k", "0.6295", "1.4 GB", "최종 제출"],
    ["2", "GPTQ W8A8 + SparseGPT 10%", "0.6246", "1.3 GB", "비정형 스파시티"],
    ["3", "GPTQ W8A8 (cal=2048, d=0.02)", "0.6208", "1.8 GB", "Hessian 캘리브레이션"],
    ["4", "GPTQ W8A8 (d=0.05)", "0.6077", "1.8 GB", "dampening 과다"],
    ["5", "GPTQ W8A8 기본 (d=0.01)", "0.6042", "1.8 GB", "기본 설정"],
    ["6", "W4A16 GPTQ (n=512)", "0.5099", "1.3 GB", "4비트 품질 손실"],
    ["7", "GPTQ W8A8 actorder=group", "0.5085", "1.8 GB", "커널 비효율"],
    ["8", "FP8 Dynamic", "0.4936", "1.4 GB", "L4 FP8 커널 미흡"],
    ["9", "W4A16 + SFT", "0.4825", "1.3 GB", "SFT 효과 미미"],
    ["10", "W8A8 Selective (5L BF16)", "0.4718", "1.8 GB", "혼합 정밀도 페널티"],
    ["11", "SFT + W4A16", "0.3202", "1.3 GB", "SFT 성능 저하"],
    ["12", "3.5세대 KD + W4A16", "0.2632", "1.3 GB", "세대 불일치"],
], col_widths=[0.7, 4.2, 1.5, 1.2, 4.7])

# 1등 행 강조
for j in range(5):
    cell = result_table.cell(1, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = HIGHLIGHT
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = BLUE

# 하단 요약
add_text(slide, 0.5, 6.4, 12, 0.5,
         "W8A8: PerfNorm≈0.99, SpeedNorm≈0.12  |  W4A16: PerfNorm≈0.50~0.57  |  FP8: SpeedNorm≈0.002",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ===========================================================
# Slide 5: 실패한 방법들
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 9.5, 0.6,
         "시도했으나 실패한 방법들", size=30, bold=True, color=BLUE)
add_text(slide, 9.5, 0.38, 3.5, 0.4,
         "총 110+ 모델 실험, 15+ 서버 제출",
         size=14, bold=True, color=RED, align=PP_ALIGN.RIGHT)
add_divider(slide)

add_table(slide, 0.3, 1.2, 12.7, 5.8, [
    ["카테고리", "시도한 방법", "결과", "실패 원인"],
    ["4비트 양자화", "W4A16 GPTQ (256/512/1024 샘플)", "Score 0.47~0.51", "1.2B에 INT4 품질 손실 과도"],
    ["FP8 양자화", "FP8 Dynamic / Static", "Score ~0.49", "L4 FP8 커널 최적화 부족"],
    ["SFT (미세조정)", "LoRA SFT → 양자화", "Score 0.32~0.48", "MANTA SFT가 1.2B 성능 저하"],
    ["Post-quant SFT", "양자화 후 norm만 SFT", "변화 없음", "학습 파라미터 <5% 불충분"],
    ["레이어 프루닝", "30L → 28L/26L + KD", "Score 0.47", "속도 이득 대비 품질 손실 과도"],
    ["FFN 프루닝", "25% FFN 폭 축소", "perf 0.273", "너무 공격적, 복구 불가"],
    ["Vocab 프루닝", "102k → 80k 토큰", "Score 0.48", "byte-fallback으로 토큰 폭증"],
    ["지식증류 (KD)", "32B/7.8B teacher → 1.2B", "Score 0.26", "가중치 변경 → 양자화 에러 증가"],
    ["SmoothQuant", "활성화 분포 평탄화", "실행 불가", "EXAONE 아키텍처 미지원"],
    ["2:4 Sparsity", "Wanda 구조적 스파시티", "실행 불가", "compute cap >= 90 필요 (L4=89)"],
    ["Selective Quant", "후반 5개 레이어 BF16 유지", "Score 0.47", "혼합 정밀도 → INT8 커널 이탈"],
    ["GPTQ Advanced", "actorder=group, block=128", "Score 0.51", "그룹 양자화 CUTLASS 비효율"],
    ["FP8 KV Cache", "KV Cache FP8 양자화", "변화 없음", "vLLM이 설정 무시"],
], col_widths=[1.8, 3.8, 2.0, 5.1])

# ===========================================================
# Slide 6: 핵심 인사이트
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 12, 0.6,
         "핵심 인사이트 4가지", size=30, bold=True, color=BLUE)
add_divider(slide)

insights = [
    ("1. L4 GPU에서 W8A8 (INT8)만 유효", [
        "• CutlassScaledMM 커널로 INT8 연산 가속 (~24%)",
        "• W8A16은 AllSpark 커널 → 느림",
        "• FP8은 vLLM 커널 최적화 부족",
        "→ 하드웨어별 최적 커널 확인 필수",
    ]),
    ("2. 단순한 양자화가 최고 성능", [
        "• QuantMod(min-max) > GPTQ(Hessian)",
        "• 복잡한 캘리브레이션이 과적합 유발",
        "• Data-Free가 오히려 일반화 우수",
        "→ vLLM INT8 CUTLASS 커널 경로 유지가 핵심",
    ]),
    ("3. 로컬 vs 서버 결과 괴리", [
        "• RTX 5060Ti와 L4: 최적 양자화 방식 다름",
        "• 양자화에 따라 vLLM 커널 선택 달라짐",
        "• 로컬에서 비슷해도 서버에서 큰 차이",
        "→ 반드시 서버 제출 기준으로 최종 판단",
    ]),
    ("4. 소형 모델(1.2B)의 양자화 특성", [
        "• INT4는 7B+에서 효과적, 1.2B는 품질 손실 과도",
        "• KD 가중치 변경 시 양자화 에러 급증 (최대 94.65)",
        "• SFT도 소형 모델에서 오히려 성능 저하",
        "→ 소형 모델은 최소한의 변형(W8A8)이 최적",
    ]),
]

for i, (title, bullets) in enumerate(insights):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.3 + row * 2.9

    add_card(slide, x, y, 5.8, 2.6)

    add_text(slide, x + 0.2, y + 0.15, 5.4, 0.4,
             title, size=16, bold=True, color=BLUE)

    tf = add_text(slide, x + 0.2, y + 0.6, 5.4, 1.8, "", size=12)
    for b in bullets:
        if b.startswith("→"):
            add_para(tf, b, size=12, bold=True, color=RED)
        else:
            add_para(tf, b, size=12, color=TEXT)

# ===========================================================
# Slide 7: 모델 크기 비교 + 재현 방법
# ===========================================================
slide = prs.slides.add_slide(blank)
add_text(slide, 0.5, 0.3, 12, 0.6,
         "모델 크기 비교 + 재현 방법", size=30, bold=True, color=BLUE)
add_divider(slide)

# 왼쪽: 비교 테이블
add_text(slide, 0.5, 1.2, 6, 0.4,
         "성능 비교", size=18, bold=True, color=BLUE)

comp_table = add_table(slide, 0.5, 1.7, 6, 3.5, [
    ["항목", "베이스 모델", "경량화 모델"],
    ["정밀도", "BF16 (16비트)", "INT8 (W8A8)"],
    ["모델 크기", "2.4 GB", "1.4 GB (42% 감소)"],
    ["max_position", "65,536", "16,384"],
    ["추론 속도 (L4)", "1,811 tok/s", "2,435 tok/s"],
    ["PerfNorm", "1.000", "0.990"],
    ["SpeedNorm", "0.000", "0.121"],
    ["Score", "0.500", "0.6295"],
], col_widths=[2, 2, 2])

# Score 행 강조
for j in range(3):
    cell = comp_table.cell(7, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = HIGHLIGHT
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = BLUE

# 오른쪽: 재현 방법
add_text(slide, 7, 1.2, 5.5, 0.4,
         "재현 방법", size=18, bold=True, color=BLUE)

# 코드 블록 배경
add_card(slide, 7, 1.7, 5.8, 4.5, RGBColor(0x2D, 0x2D, 0x2D))

code_tf = add_text(slide, 7.2, 1.8, 5.4, 4.3, "", size=11,
                   color=GREEN, font_name="Consolas")
code_lines = [
    ("# 1. 환경 설정", True, RGBColor(0x6A, 0x99, 0x55)),
    ("conda create -n lgaimers python=3.10", False, WHITE),
    ("conda activate lgaimers", False, WHITE),
    ("pip install -r requirements.txt", False, WHITE),
    ("", False, WHITE),
    ("# 2. 경량화 실행 (~10초)", True, RGBColor(0x6A, 0x99, 0x55)),
    ("python reproduce.py", False, SKY),
    ("", False, WHITE),
    ("# 핵심 코드 (reproduce.py):", True, RGBColor(0x6A, 0x99, 0x55)),
    ("QuantizationModifier(", False, RGBColor(0xDC, 0xDC, 0xAA)),
    ('    scheme="W8A8",', False, RGBColor(0xCE, 0x91, 0x78)),
    ('    targets=["Linear"],', False, RGBColor(0xCE, 0x91, 0x78)),
    ('    ignore=["embed_tokens",', False, RGBColor(0xCE, 0x91, 0x78)),
    ('           "lm_head"]', False, RGBColor(0xCE, 0x91, 0x78)),
    (")", False, RGBColor(0xDC, 0xDC, 0xAA)),
    ("", False, WHITE),
    ("# config.json 수정:", True, RGBColor(0x6A, 0x99, 0x55)),
    ("# max_position: 65536 → 16384", False, RGBColor(0x6A, 0x99, 0x55)),
]
for text, is_bold, clr in code_lines:
    add_para(code_tf, text, size=11, bold=is_bold, color=clr, font_name="Consolas")

# 하단 요약 바
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7)
)
shape.fill.solid()
shape.fill.fore_color.rgb = HIGHLIGHT
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "   GPU: 12GB+  |  양자화: ~10초  |  외부 데이터: MANTA-1M (LG AI Research 공개)  |  결과: submit.zip (1.2GB)"
p.font.size = Pt(13)
p.font.color.rgb = BLUE
p.font.name = "맑은 고딕"
p.alignment = PP_ALIGN.CENTER

# ===========================================================
# 저장
# ===========================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "solution.pptx")
prs.save(out_path)
print(f"PPT 저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}")
